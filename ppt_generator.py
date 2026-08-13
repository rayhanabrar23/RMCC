"""
ppt_generator.py
------------------
Membangun file .pptx dari KomiteSummary (lihat excel_parser.py).
Pakai python-pptx (bukan pptxgenjs) supaya jalan murni di Python —
cocok untuk deploy di Streamlit Community Cloud tanpa Node.js.

Chart digambar pakai matplotlib lalu ditempel sebagai gambar,
karena python-pptx tidak punya API chart senyaman pptxgenjs.
"""

import io

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ─────────────────────────────────────────────
# PALET WARNA & KONSTAN
# ─────────────────────────────────────────────
NAVY = RGBColor(0x1C, 0x2B, 0x4A)
NAVY_LIGHT = RGBColor(0x2E, 0x42, 0x66)
ACCENT = RGBColor(0xC8, 0x96, 0x3E)
GREEN = RGBColor(0x1E, 0x8E, 0x5A)
RED = RGBColor(0xC0, 0x39, 0x2B)
GRAY_TEXT = RGBColor(0x5A, 0x64, 0x72)
BG_LIGHT = RGBColor(0xF4, 0xF6, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# hex tanpa '#' — dipakai khusus untuk matplotlib
NAVY_HEX = "#1C2B4A"
NAVY_LIGHT_HEX = "#2E4266"
GRAY_HEX = "#8A94A3"

FONT = "Calibri"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _fmt_rp_triliun(v):
    return f"{v / 1e12:.2f} T"


def _fmt_rp_miliar(v):
    return f"Rp {v / 1e9:,.1f} Miliar".replace(",", ".")


def _fmt_pct(v):
    return f"{v * 100:.1f}%".replace(".", ",")


# ─────────────────────────────────────────────
# HELPERS SHAPE python-pptx
# ─────────────────────────────────────────────
def _set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _no_shadow(shape):
    shape.shadow.inherit = False


def _add_shadow(shape, alpha=0.15, blur_pt=8, dist_pt=2):
    """Soft drop shadow sederhana (arah bawah) via XML langsung."""
    sp = shape._element.spPr
    effect_lst = sp.makeelement(qn("a:effectLst"), {})
    outer_shdw = sp.makeelement(qn("a:outerShdw"), {
        "blurRad": str(Pt(blur_pt)),
        "dist": str(Pt(dist_pt)),
        "dir": "5400000",
        "rotWithShape": "0",
    })
    clr = sp.makeelement(qn("a:srgbClr"), {"val": "1C2B4A"})
    alpha_el = sp.makeelement(qn("a:alpha"), {"val": str(int(alpha * 100000))})
    clr.append(alpha_el)
    outer_shdw.append(clr)
    effect_lst.append(outer_shdw)
    sp.append(effect_lst)


def _textbox(slide, x, y, w, h, text, size=12, bold=False, italic=False,
             color=NAVY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT, line_spacing=1.0, char_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing != 1.0:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color
    return tb


def _rich_textbox(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """runs: list of dict(text, size, bold, color)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for r in runs:
        run = p.add_run()
        run.text = r["text"]
        run.font.size = Pt(r.get("size", 12))
        run.font.bold = r.get("bold", False)
        run.font.name = FONT
        run.font.color.rgb = r.get("color", NAVY)
    return tb


def _rounded_rect(slide, x, y, w, h, fill_color, radius=0.06, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = radius
    _set_fill(shp, fill_color)
    if shadow:
        _add_shadow(shp)
    else:
        _no_shadow(shp)
    return shp


def _footer(slide, page_num, periode_label):
    _textbox(slide, Inches(0.5), Inches(7.12), Inches(9.0), Inches(0.3),
             f"Komite Haircut & Concentration Limit — {periode_label}  |  Internal / Confidential",
             size=9, color=RGBColor(0x9A, 0xA3, 0xAF))
    _textbox(slide, Inches(12.5), Inches(7.12), Inches(0.5), Inches(0.3),
             str(page_num), size=9, color=RGBColor(0x9A, 0xA3, 0xAF), align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────
# CHART (matplotlib -> PNG in-memory)
# ─────────────────────────────────────────────
def _chart_group_distribution(group_dist):
    fig, ax = plt.subplots(figsize=(8.0, 4.7), dpi=200)
    names = [f"{g['name']}\n({g['n_stocks']} saham)" for g in group_dist]
    values = [g["value"] / 1e12 for g in group_dist]
    bars = ax.bar(names, values, color=NAVY_HEX, width=0.55)
    for b, v in zip(bars, values):
        ax.text(b.get_x() + b.get_width() / 2, v + max(values) * 0.02, f"{v:.2f} T",
                ha="center", va="bottom", fontsize=11, color=NAVY_HEX, fontweight="bold")
    ax.set_ylabel("Rp Triliun", fontsize=10, color=GRAY_HEX)
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.tick_params(axis="x", labelsize=11, colors=NAVY_HEX)
    ax.tick_params(axis="y", labelsize=9, colors=GRAY_HEX)
    ax.yaxis.grid(True, color="#E4E8EE", linewidth=1)
    ax.set_axisbelow(True)
    ax.set_ylim(0, max(values) * 1.2)
    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf


def _chart_reason_counts(reason_counts, top_n=6):
    """reason_counts: dict alasan -> jumlah. Digroup jadi kategori pendek."""
    items = sorted(reason_counts.items(), key=lambda kv: -kv[1])
    fig, ax = plt.subplots(figsize=(7.2, 4.7), dpi=200)
    labels = [k if len(k) <= 28 else k[:26] + "…" for k, _ in items[:top_n]]
    values = [v for _, v in items[:top_n]]
    y_pos = range(len(labels))
    bars = ax.barh(list(y_pos), values, color=NAVY_LIGHT_HEX, height=0.55)
    for b, v in zip(bars, values):
        ax.text(v + max(values) * 0.02, b.get_y() + b.get_height() / 2, str(v),
                va="center", fontsize=11, color=NAVY_HEX, fontweight="bold")
    ax.set_yticks(list(y_pos))
    ax.set_yticklabels(labels, fontsize=10.5, color=NAVY_HEX)
    ax.invert_yaxis()
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.tick_params(axis="x", labelsize=9, colors=GRAY_HEX)
    ax.xaxis.grid(True, color="#E4E8EE", linewidth=1)
    ax.set_axisbelow(True)
    ax.set_xlim(0, max(values) * 1.25)
    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf


# ─────────────────────────────────────────────
# SLIDE BUILDERS
# ─────────────────────────────────────────────
def _slide_title(prs, summary, company_name):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_fill(bg, NAVY)
    _no_shadow(bg)

    # dekorasi lingkaran (bukan stripe/garis aksen)
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.6), Inches(-1.2), Inches(4.2), Inches(4.2))
    _set_fill(c1, NAVY_LIGHT)
    _no_shadow(c1)
    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.6), Inches(5.2), Inches(2.6), Inches(2.6))
    c2.fill.solid()
    c2.fill.fore_color.rgb = ACCENT
    c2.fill.transparency = 0.8 if hasattr(c2.fill, "transparency") else None
    c2.line.fill.background()
    _no_shadow(c2)

    _textbox(slide, Inches(0.9), Inches(2.55), Inches(11.5), Inches(0.9),
              "KOMITE HAIRCUT & CONCENTRATION LIMIT", size=32, bold=True, color=WHITE)
    _textbox(slide, Inches(0.9), Inches(3.35), Inches(11.5), Inches(0.5),
              f"Ringkasan Hasil Pembahasan — Periode {summary.periode_label}", size=17,
              color=RGBColor(0xC9, 0xD2, 0xDE))
    _textbox(slide, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.4),
              f"Risk Management & Control Division  •  {company_name}",
              size=12, color=RGBColor(0x8F, 0xA0, 0xBA))
    return slide


def _metric_card(slide, x, y, w, h, value, label, sub, color):
    card = _rounded_rect(slide, x, y, w, h, BG_LIGHT, radius=0.08, shadow=True)
    _textbox(slide, x + Inches(0.15), y + Inches(0.22), w - Inches(0.3), Inches(0.75),
              value, size=32, bold=True, color=color)
    _textbox(slide, x + Inches(0.15), y + Inches(0.98), w - Inches(0.3), Inches(0.35),
              label, size=10.5, bold=True, color=NAVY)
    _textbox(slide, x + Inches(0.15), y + Inches(1.34), w - Inches(0.3), Inches(0.4),
              sub, size=10.5, color=GRAY_TEXT)


def _slide_summary(prs, summary, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_fill(bg, WHITE)
    _no_shadow(bg)

    _textbox(slide, Inches(0.6), Inches(0.4), Inches(10), Inches(0.6),
              "Ringkasan Eksekutif", size=26, bold=True, color=NAVY)
    _textbox(slide, Inches(0.6), Inches(0.95), Inches(10), Inches(0.35),
              "Snapshot data hasil pembahasan komite", size=13, color=GRAY_TEXT)

    cards = [
        (str(summary.total_saham), "TOTAL SAHAM DIANALISIS", "efek margin/repo", NAVY),
        (str(summary.haircut_naik + summary.haircut_turun), "HAIRCUT BERUBAH",
         f"{summary.haircut_naik} naik  •  {summary.haircut_turun} turun", ACCENT),
        (str(summary.uma_count), "TERKENA UMA", "saham", RED),
        (str(summary.saham_baru_count), "SAHAM BARU", f"periode {summary.saham_baru_periode_label}", GREEN),
    ]
    card_w, gap, start_x, y, h = Inches(2.75), Inches(0.35), Inches(0.6), Inches(1.65), Inches(1.9)
    for i, (val, label, sub, color) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        _metric_card(slide, x, y, card_w, h, val, label, sub, color)

    box = _rounded_rect(slide, Inches(0.6), Inches(3.9), Inches(12.13), Inches(1.55), NAVY, radius=0.08)
    _textbox(slide, Inches(1.0), Inches(4.1), Inches(5), Inches(0.35),
              "PARAMETER PERHITUNGAN", size=11, bold=True, color=ACCENT)
    _rich_textbox(slide, Inches(1.0), Inches(4.5), Inches(11.3), Inches(0.5), [
        {"text": "PEI Equity: ", "bold": True, "color": WHITE, "size": 15},
        {"text": f"{_fmt_rp_miliar(summary.pei_equity)}     ", "color": RGBColor(0xC9, 0xD2, 0xDE), "size": 15},
        {"text": "Max Financing Ratio: ", "bold": True, "color": WHITE, "size": 15},
        {"text": _fmt_pct(summary.max_fin_ratio), "color": RGBColor(0xC9, 0xD2, 0xDE), "size": 15},
    ])
    _textbox(slide, Inches(1.0), Inches(5.0), Inches(11.3), Inches(0.35),
              "Parameter dasar penentuan Concentration Limit.", size=10.5, italic=True,
              color=RGBColor(0x9A, 0xA3, 0xAF))

    top_reason = max(summary.reason_counts.items(), key=lambda kv: kv[1])[0] if summary.reason_counts else "-"
    _textbox(slide, Inches(0.6), Inches(5.75), Inches(12.13), Inches(0.9),
              f"Poin utama: {summary.haircut_naik + summary.haircut_turun} saham mengalami penyesuaian haircut "
              f"pada periode ini. Alasan paling dominan: \"{top_reason}\".",
              size=13, color=NAVY)

    _footer(slide, page_num, summary.periode_label)
    return slide


def _slide_group_distribution(prs, summary, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_fill(bg, WHITE)
    _no_shadow(bg)

    _textbox(slide, Inches(0.6), Inches(0.4), Inches(10), Inches(0.6),
              "Distribusi Konsentrasi Efek Jaminan", size=26, bold=True, color=NAVY)
    _textbox(slide, Inches(0.6), Inches(0.95), Inches(11), Inches(0.35),
              "Berdasarkan kelompok haircut (Group HC) — Max Collateral Value setelah haircut",
              size=13, color=GRAY_TEXT)

    if summary.group_dist:
        chart_buf = _chart_group_distribution(summary.group_dist)
        slide.shapes.add_picture(chart_buf, Inches(0.5), Inches(1.5), width=Inches(8.1))

        panel = _rounded_rect(slide, Inches(9.0), Inches(1.5), Inches(3.75), Inches(4.9), BG_LIGHT, radius=0.08)
        _textbox(slide, Inches(9.3), Inches(1.7), Inches(3.2), Inches(0.35),
                  "% TERHADAP TOTAL", size=11, bold=True, color=NAVY)
        colors = [NAVY, NAVY_LIGHT, ACCENT, RED]
        for i, g in enumerate(summary.group_dist[:4]):
            y = Inches(2.25) + Emu(int(Inches(1.0)) * i)
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.3), y, Inches(0.14), Inches(0.55))
            _set_fill(bar, colors[i % len(colors)])
            _no_shadow(bar)
            _textbox(slide, Inches(9.55), y - Inches(0.03), Inches(2.2), Inches(0.35),
                      g["name"], size=13, bold=True, color=NAVY)
            _textbox(slide, Inches(9.55), y + Inches(0.28), Inches(2.2), Inches(0.35),
                      _fmt_pct(g["pct"]), size=15, bold=True, color=colors[i % len(colors)])

        total_n = sum(g["n_stocks"] for g in summary.group_dist)
        _textbox(slide, Inches(9.3), Inches(5.9), Inches(3.3), Inches(0.4),
                  f"Total: {total_n} saham  |  {_fmt_rp_triliun(summary.total_collateral_value)}",
                  size=10.5, italic=True, color=GRAY_TEXT)
    else:
        _textbox(slide, Inches(0.6), Inches(2.5), Inches(11), Inches(1),
                  "Data distribusi Group HC tidak tersedia di file sumber.", size=14, color=RED)

    _footer(slide, page_num, summary.periode_label)
    return slide


def _slide_reason_breakdown(prs, summary, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_fill(bg, WHITE)
    _no_shadow(bg)

    total_changed = summary.haircut_naik + summary.haircut_turun
    _textbox(slide, Inches(0.6), Inches(0.4), Inches(11), Inches(0.6),
              f"Alasan Perubahan Haircut — {total_changed} Saham", size=25, bold=True, color=NAVY)
    _textbox(slide, Inches(0.6), Inches(0.95), Inches(10), Inches(0.35),
              "Dibandingkan hasil pembahasan periode sebelumnya", size=13, color=GRAY_TEXT)

    # kelompokkan alasan yang menyebut "UMA" jadi satu kategori ringkas untuk chart & panel
    uma_related = sum(v for k, v in summary.reason_counts.items() if "UMA" in k.upper())
    non_uma = {k: v for k, v in summary.reason_counts.items() if "UMA" not in k.upper()}
    grouped = dict(non_uma)
    if uma_related:
        grouped["Mengikuti Notasi UMA dari BEI"] = uma_related

    if grouped:
        chart_buf = _chart_reason_counts(grouped)
        slide.shapes.add_picture(chart_buf, Inches(0.5), Inches(1.5), width=Inches(7.4))
    else:
        _textbox(slide, Inches(0.6), Inches(2.5), Inches(7), Inches(1),
                  "Tidak ada perubahan haircut pada periode ini.", size=14, color=GRAY_TEXT)

    panel = _rounded_rect(slide, Inches(8.35), Inches(1.5), Inches(4.4), Inches(4.9), NAVY, radius=0.08)
    _textbox(slide, Inches(8.7), Inches(1.75), Inches(3.7), Inches(0.35),
              "PERHATIAN KOMITE", size=11, bold=True, color=ACCENT)
    _rich_textbox(slide, Inches(8.7), Inches(2.25), Inches(3.7), Inches(1.4), [
        {"text": f"{uma_related} saham ", "bold": True, "color": WHITE, "size": 15},
        {"text": "mengalami penyesuaian haircut mengikuti notasi UMA (Unusual Market Activity) dari BEI.",
         "color": RGBColor(0xC9, 0xD2, 0xDE), "size": 12.5},
    ])
    line = slide.shapes.add_connector(1, Inches(8.7), Inches(3.75), Inches(12.4), Inches(3.75))
    line.line.color.rgb = RGBColor(0x3A, 0x4C, 0x6E)
    line.line.width = Pt(1)

    _textbox(slide, Inches(8.7), Inches(3.95), Inches(3.7), Inches(0.3),
              "TOTAL PERUBAHAN", size=10, bold=True, color=RGBColor(0x8F, 0xA0, 0xBA))
    _textbox(slide, Inches(8.7), Inches(4.25), Inches(3.7), Inches(0.4),
              f"{total_changed} saham ({summary.haircut_naik} naik, {summary.haircut_turun} turun)",
              size=13.5, bold=True, color=WHITE)

    line2 = slide.shapes.add_connector(1, Inches(8.7), Inches(4.85), Inches(12.4), Inches(4.85))
    line2.line.color.rgb = RGBColor(0x3A, 0x4C, 0x6E)
    line2.line.width = Pt(1)

    _textbox(slide, Inches(8.7), Inches(5.05), Inches(3.7), Inches(0.3),
              "TINDAK LANJUT", size=10, bold=True, color=RGBColor(0x8F, 0xA0, 0xBA))
    _textbox(slide, Inches(8.7), Inches(5.35), Inches(3.7), Inches(0.9),
              "Divisi memantau notasi UMA berjalan dan menyesuaikan haircut bila status berubah pada periode berikutnya.",
              size=11.5, color=RGBColor(0xC9, 0xD2, 0xDE))

    _footer(slide, page_num, summary.periode_label)
    return slide


def _slide_recommendations(prs, summary, page_num, custom_points=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_fill(bg, BG_LIGHT)
    _no_shadow(bg)

    _textbox(slide, Inches(0.6), Inches(0.5), Inches(10.5), Inches(0.6),
              "Rekomendasi & Tindak Lanjut", size=26, bold=True, color=NAVY)

    default_points = [
        ("Penerapan Haircut & CL Baru",
         f"Berlaku efektif untuk transaksi margin/repo mulai periode {summary.periode_label}, "
         f"mengacu pada hasil pembahasan komite."),
        ("Monitoring Saham Ter-UMA",
         f"{summary.uma_count} saham dengan status UMA perlu pemantauan lanjutan bila ada perubahan notasi dari BEI."),
        ("Review Saham Margin Baru",
         f"{summary.saham_baru_count} saham baru masuk kategori margin — perlu konfirmasi kesiapan sistem OP/collateral."),
    ]
    points = custom_points if custom_points else default_points

    y0, row_h = Inches(1.5), Inches(1.55)
    for i, (title, desc) in enumerate(points):
        y = y0 + i * row_h
        _rounded_rect(slide, Inches(0.6), y, Inches(12.13), Inches(1.3), WHITE, radius=0.08, shadow=True)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), y + Inches(0.35), Inches(0.6), Inches(0.6))
        _set_fill(circle, NAVY)
        _no_shadow(circle)
        tf = circle.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(i + 1)
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = FONT

        _textbox(slide, Inches(1.85), y + Inches(0.18), Inches(10.6), Inches(0.4), title,
                  size=15, bold=True, color=NAVY)
        _textbox(slide, Inches(1.85), y + Inches(0.58), Inches(10.6), Inches(0.6), desc,
                  size=12, color=GRAY_TEXT)

    _footer(slide, page_num, summary.periode_label)
    return slide


# ─────────────────────────────────────────────
# ENTRY POINT
# ─────────────────────────────────────────────
def build_presentation(summary, company_name="PT Pendanaan Efek Indonesia", custom_recommendations=None):
    """Return io.BytesIO berisi file .pptx siap didownload."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_title(prs, summary, company_name)
    _slide_summary(prs, summary, page_num=2)
    _slide_group_distribution(prs, summary, page_num=3)
    _slide_reason_breakdown(prs, summary, page_num=4)
    _slide_recommendations(prs, summary, page_num=5, custom_points=custom_recommendations)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
